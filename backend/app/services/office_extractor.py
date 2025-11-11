"""
Office Document Text Extraction Service.

This module provides text extraction from Microsoft Office documents including:
- Word documents (.docx, .doc)
- Excel spreadsheets (.xlsx, .xls)
- PowerPoint presentations (.pptx, .ppt)

Features:
- Modern Office formats (OOXML): docx, xlsx, pptx
- Legacy Office formats (OLE): doc, xls, ppt
- Preserves document structure where possible
- Handles tables, bullet points, and formatting
- Extracts text from embedded objects
- Character encoding detection for legacy formats

Usage:
    from app.services.office_extractor import OfficeExtractor

    extractor = OfficeExtractor()
    result = await extractor.extract_text(file_content, "document.docx")
"""

import io
import zipfile
from pathlib import Path
from typing import Any

import chardet
import olefile
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.utils.exceptions import DocumentProcessingError
from app.utils.logger import get_logger

logger = get_logger(__name__)


class OfficeExtractor:
    """Extract text from Microsoft Office documents."""

    def __init__(self):
        """Initialize Office extractor."""
        self.max_text_length = 10_000_000  # 10MB text limit

    async def extract_text(
        self,
        file_content: bytes,
        filename: str,
    ) -> dict[str, Any]:
        """
        Extract text from Office document.

        Args:
            file_content: Document file bytes
            filename: Original filename with extension

        Returns:
            Dictionary with extracted text and metadata:
            {
                "text": str,
                "method": str,
                "pages": int,
                "confidence": float,
                "metadata": dict
            }

        Raises:
            DocumentProcessingError: If extraction fails
        """
        file_extension = Path(filename).suffix.lower()

        try:
            # Word documents
            if file_extension == ".docx":
                return await self._extract_docx(file_content)
            elif file_extension == ".doc":
                return await self._extract_doc(file_content)

            # Excel spreadsheets
            elif file_extension == ".xlsx":
                return await self._extract_xlsx(file_content)
            elif file_extension == ".xls":
                return await self._extract_xls(file_content)

            # PowerPoint presentations
            elif file_extension == ".pptx":
                return await self._extract_pptx(file_content)
            elif file_extension == ".ppt":
                return await self._extract_ppt(file_content)

            else:
                raise DocumentProcessingError(
                    message=f"Unsupported Office format: {file_extension}",
                    details={"filename": filename, "extension": file_extension},
                )

        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"Office extraction failed for {filename}: {e}", exc_info=True)
            raise DocumentProcessingError(
                message=f"Failed to extract text from Office document: {str(e)}",
                details={"filename": filename, "error": str(e)},
            )

    # =========================================================================
    # Word Document Extraction
    # =========================================================================

    async def _extract_docx(self, file_content: bytes) -> dict[str, Any]:
        """
        Extract text from modern Word document (.docx).

        Args:
            file_content: File bytes

        Returns:
            Extraction result with text and metadata
        """
        try:
            doc = Document(io.BytesIO(file_content))

            text_parts = []
            paragraph_count = 0
            table_count = 0

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
                    paragraph_count += 1

            # Extract tables
            for table in doc.tables:
                table_count += 1
                table_text = self._extract_table_text(table)
                if table_text:
                    text_parts.append(f"\n[TABLE {table_count}]\n{table_text}\n")

            text = "\n".join(text_parts)

            # Extract metadata
            core_properties = doc.core_properties
            metadata = {
                "author": core_properties.author or "",
                "title": core_properties.title or "",
                "created": str(core_properties.created) if core_properties.created else "",
                "modified": str(core_properties.modified) if core_properties.modified else "",
                "paragraph_count": paragraph_count,
                "table_count": table_count,
            }

            return {
                "text": text,
                "method": "python-docx",
                "pages": len(doc.sections),  # Approximate page count
                "confidence": 100.0,
                "metadata": metadata,
            }

        except Exception as e:
            raise DocumentProcessingError(
                message=f"Failed to parse DOCX file: {str(e)}",
                details={"error": str(e)},
            )

    def _extract_table_text(self, table) -> str:
        """Extract text from a Word table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):  # Only include non-empty rows
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    async def _extract_doc(self, file_content: bytes) -> dict[str, Any]:
        """
        Extract text from legacy Word document (.doc).

        Args:
            file_content: File bytes

        Returns:
            Extraction result with text and metadata
        """
        try:
            # Check if it's a valid OLE file
            if not olefile.isOleFile(io.BytesIO(file_content)):
                raise DocumentProcessingError(
                    message="File is not a valid legacy Word document",
                    details={"format": "OLE"},
                )

            ole = olefile.OleFileIO(io.BytesIO(file_content))

            # Try to extract text from WordDocument stream
            text_parts = []

            # Method 1: Try WordDocument stream
            if ole.exists("WordDocument"):
                word_stream = ole.openstream("WordDocument")
                raw_data = word_stream.read()

                # Detect encoding
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result["encoding"] or "utf-8"
                confidence = encoding_result["confidence"]

                try:
                    # Extract printable text
                    text = raw_data.decode(encoding, errors="ignore")
                    # Clean up binary artifacts
                    text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                    text_parts.append(text)
                except Exception as decode_error:
                    logger.warning(f"Failed to decode .doc with {encoding}: {decode_error}")

            # Method 2: Try 1Table stream (contains text)
            if ole.exists("1Table"):
                try:
                    table_stream = ole.openstream("1Table")
                    table_data = table_stream.read()
                    text = table_data.decode("utf-16-le", errors="ignore")
                    text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                    if text.strip():
                        text_parts.append(text)
                except Exception as e:
                    logger.warning(f"Failed to extract from 1Table: {e}")

            ole.close()

            if not text_parts:
                # Fallback: Try to extract any readable text
                text = file_content.decode("latin-1", errors="ignore")
                text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                text_parts.append(text)

            final_text = "\n".join(text_parts)

            # Basic cleanup
            final_text = "\n".join(line.strip() for line in final_text.split("\n") if line.strip())

            return {
                "text": final_text,
                "method": "olefile",
                "pages": 1,  # Cannot determine pages from .doc easily
                "confidence": 70.0,  # Lower confidence for legacy format
                "metadata": {
                    "format": "legacy_word",
                    "note": "Legacy format - formatting may be lost",
                },
            }

        except Exception as e:
            raise DocumentProcessingError(
                message=f"Failed to parse DOC file: {str(e)}",
                details={"error": str(e)},
            )

    # =========================================================================
    # Excel Spreadsheet Extraction
    # =========================================================================

    async def _extract_xlsx(self, file_content: bytes) -> dict[str, Any]:
        """
        Extract text from modern Excel spreadsheet (.xlsx).

        Args:
            file_content: File bytes

        Returns:
            Extraction result with text and metadata
        """
        try:
            workbook = load_workbook(io.BytesIO(file_content), data_only=True)

            text_parts = []
            sheet_count = len(workbook.sheetnames)
            total_rows = 0
            total_cells = 0

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                # Add sheet header
                text_parts.append(f"\n[SHEET: {sheet_name}]\n")

                # Extract cell values
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty rows
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(val.strip() for val in row_values):
                        rows.append(" | ".join(row_values))
                        total_rows += 1
                        total_cells += len(row_values)

                if rows:
                    text_parts.append("\n".join(rows))

            text = "\n".join(text_parts)

            metadata = {
                "sheet_count": sheet_count,
                "total_rows": total_rows,
                "total_cells": total_cells,
                "sheets": workbook.sheetnames,
            }

            workbook.close()

            return {
                "text": text,
                "method": "openpyxl",
                "pages": sheet_count,  # Each sheet is like a page
                "confidence": 100.0,
                "metadata": metadata,
            }

        except Exception as e:
            raise DocumentProcessingError(
                message=f"Failed to parse XLSX file: {str(e)}",
                details={"error": str(e)},
            )

    async def _extract_xls(self, file_content: bytes) -> dict[str, Any]:
        """
        Extract text from legacy Excel spreadsheet (.xls).

        Args:
            file_content: File bytes

        Returns:
            Extraction result with text and metadata
        """
        try:
            # Check if it's a valid OLE file
            if not olefile.isOleFile(io.BytesIO(file_content)):
                raise DocumentProcessingError(
                    message="File is not a valid legacy Excel document",
                    details={"format": "OLE"},
                )

            ole = olefile.OleFileIO(io.BytesIO(file_content))

            text_parts = []

            # Try to extract text from Workbook stream
            if ole.exists("Workbook"):
                workbook_stream = ole.openstream("Workbook")
                raw_data = workbook_stream.read()

                # Detect encoding
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result["encoding"] or "utf-8"

                try:
                    # Extract printable text
                    text = raw_data.decode(encoding, errors="ignore")
                    # Clean up binary artifacts
                    text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                    text_parts.append(text)
                except Exception as decode_error:
                    logger.warning(f"Failed to decode .xls with {encoding}: {decode_error}")

            ole.close()

            if not text_parts:
                # Fallback: Try to extract any readable text
                text = file_content.decode("latin-1", errors="ignore")
                text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                text_parts.append(text)

            final_text = "\n".join(text_parts)

            # Basic cleanup
            final_text = "\n".join(line.strip() for line in final_text.split("\n") if line.strip())

            return {
                "text": final_text,
                "method": "olefile",
                "pages": 1,  # Cannot determine sheets from .xls easily
                "confidence": 60.0,  # Lower confidence for legacy format
                "metadata": {
                    "format": "legacy_excel",
                    "note": "Legacy format - structure may be lost",
                },
            }

        except Exception as e:
            raise DocumentProcessingError(
                message=f"Failed to parse XLS file: {str(e)}",
                details={"error": str(e)},
            )

    # =========================================================================
    # PowerPoint Presentation Extraction
    # =========================================================================

    async def _extract_pptx(self, file_content: bytes) -> dict[str, Any]:
        """
        Extract text from modern PowerPoint presentation (.pptx).

        Args:
            file_content: File bytes

        Returns:
            Extraction result with text and metadata
        """
        try:
            presentation = Presentation(io.BytesIO(file_content))

            text_parts = []
            slide_count = len(presentation.slides)

            for slide_num, slide in enumerate(presentation.slides, start=1):
                # Add slide header
                text_parts.append(f"\n[SLIDE {slide_num}]\n")

                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                    # Extract text from tables
                    if shape.has_table:
                        table_text = self._extract_pptx_table_text(shape.table)
                        if table_text:
                            slide_text.append(f"\n[TABLE]\n{table_text}\n")

                if slide_text:
                    text_parts.append("\n".join(slide_text))

                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text_parts.append(f"\n[NOTES]\n{notes_text}\n")

            text = "\n".join(text_parts)

            metadata = {
                "slide_count": slide_count,
            }

            return {
                "text": text,
                "method": "python-pptx",
                "pages": slide_count,
                "confidence": 100.0,
                "metadata": metadata,
            }

        except Exception as e:
            raise DocumentProcessingError(
                message=f"Failed to parse PPTX file: {str(e)}",
                details={"error": str(e)},
            )

    def _extract_pptx_table_text(self, table) -> str:
        """Extract text from a PowerPoint table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):  # Only include non-empty rows
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    async def _extract_ppt(self, file_content: bytes) -> dict[str, Any]:
        """
        Extract text from legacy PowerPoint presentation (.ppt).

        Args:
            file_content: File bytes

        Returns:
            Extraction result with text and metadata
        """
        try:
            # Check if it's a valid OLE file
            if not olefile.isOleFile(io.BytesIO(file_content)):
                raise DocumentProcessingError(
                    message="File is not a valid legacy PowerPoint document",
                    details={"format": "OLE"},
                )

            ole = olefile.OleFileIO(io.BytesIO(file_content))

            text_parts = []

            # Try to extract text from PowerPoint Document stream
            if ole.exists("PowerPoint Document"):
                ppt_stream = ole.openstream("PowerPoint Document")
                raw_data = ppt_stream.read()

                # Detect encoding
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result["encoding"] or "utf-8"

                try:
                    # Extract printable text
                    text = raw_data.decode(encoding, errors="ignore")
                    # Clean up binary artifacts
                    text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                    text_parts.append(text)
                except Exception as decode_error:
                    logger.warning(f"Failed to decode .ppt with {encoding}: {decode_error}")

            # Try Current User stream
            if ole.exists("Current User"):
                try:
                    user_stream = ole.openstream("Current User")
                    user_data = user_stream.read()
                    text = user_data.decode("utf-16-le", errors="ignore")
                    text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                    if text.strip():
                        text_parts.append(text)
                except Exception as e:
                    logger.warning(f"Failed to extract from Current User: {e}")

            ole.close()

            if not text_parts:
                # Fallback: Try to extract any readable text
                text = file_content.decode("latin-1", errors="ignore")
                text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
                text_parts.append(text)

            final_text = "\n".join(text_parts)

            # Basic cleanup
            final_text = "\n".join(line.strip() for line in final_text.split("\n") if line.strip())

            return {
                "text": final_text,
                "method": "olefile",
                "pages": 1,  # Cannot determine slides from .ppt easily
                "confidence": 60.0,  # Lower confidence for legacy format
                "metadata": {
                    "format": "legacy_powerpoint",
                    "note": "Legacy format - structure may be lost",
                },
            }

        except Exception as e:
            raise DocumentProcessingError(
                message=f"Failed to parse PPT file: {str(e)}",
                details={"error": str(e)},
            )
